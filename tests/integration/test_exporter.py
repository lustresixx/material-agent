from __future__ import annotations

from contextlib import suppress
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from localdeck.rendering.exporter import HTMLExporter
from localdeck.rendering.verifier import PPTXVerifier

FIXTURE = Path(__file__).parents[1] / "fixtures" / "slides" / "valid.html"
SLIDE_FIXTURES = FIXTURE.parent


async def test_exporter_creates_editable_two_slide_pptx(tmp_path: Path) -> None:
    slides_dir = tmp_path / "slides"
    slides_dir.mkdir()
    html = FIXTURE.read_text(encoding="utf-8")
    (slides_dir / "slide_01.html").write_text(html, encoding="utf-8")
    (slides_dir / "slide_02.html").write_text(
        html.replace("LocalDeck", "第二页"), encoding="utf-8"
    )
    output = tmp_path / "deck.pptx"

    await HTMLExporter().export(slides_dir, output, "16:9")
    verification = PPTXVerifier().verify(output, expected_slides=2)

    presentation = Presentation(output)
    text_shapes = [
        shape
        for slide in presentation.slides
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False) and shape.text.strip()
    ]
    assert verification.slide_count == 2
    assert verification.text_shape_count >= 2
    assert text_shapes


async def test_exporter_does_not_replace_existing_output_on_failure(
    tmp_path: Path,
) -> None:
    slides_dir = tmp_path / "slides"
    slides_dir.mkdir()
    (slides_dir / "slide_01.html").write_text("<broken", encoding="utf-8")
    output = tmp_path / "deck.pptx"
    output.write_bytes(b"existing")

    with suppress(Exception):
        await HTMLExporter().export(slides_dir, output, "16:9")

    assert output.read_bytes() == b"existing"


async def test_exporter_preserves_gradient_panel_and_svg(tmp_path: Path) -> None:
    presentation = await _export_fixture(tmp_path, "gradient.html")
    slide = presentation.slides[0]
    pictures = [
        shape for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
    ]
    title = next(shape for shape in slide.shapes if "数智化合作" in shape.text)

    image_relationships = [
        relationship
        for relationship in slide.part.rels.values()
        if relationship.reltype.endswith("/image")
    ]

    assert pictures
    assert len(image_relationships) >= 2
    assert title.text_frame.paragraphs[0].runs[0].font.color.rgb is not None
    assert str(title.text_frame.paragraphs[0].runs[0].font.color.rgb) == "FFFFFF"


async def test_exporter_preserves_independent_side_borders(tmp_path: Path) -> None:
    presentation = await _export_fixture(tmp_path, "per_side_border.html")
    slide = presentation.slides[0]
    line_colors = {
        str(shape.line.color.rgb)
        for shape in slide.shapes
        if shape.line.color.type is not None and shape.line.color.rgb is not None
    }

    assert "E60012" in line_colors
    assert "AEB8CA" in line_colors


async def test_exporter_preserves_rich_source_text_and_rounded_image(
    tmp_path: Path,
) -> None:
    presentation = await _export_fixture(tmp_path, "rich_text.html")
    slide = presentation.slides[0]
    source = next(
        shape
        for shape in slide.shapes
        if shape.has_text_frame and "来源\uFF1A" in shape.text
    )
    runs = source.text_frame.paragraphs[0].runs
    pictures = [
        shape for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
    ]

    assert source.text == "来源\uFF1A华为官方公开资料"
    assert len(runs) >= 2
    assert runs[0].font.bold is True
    assert any(run.font.bold is not True for run in runs[1:])
    assert pictures


async def _export_fixture(tmp_path: Path, fixture_name: str) -> Presentation:
    slides_dir = tmp_path / fixture_name.removesuffix(".html")
    slides_dir.mkdir()
    fixture = SLIDE_FIXTURES / fixture_name
    (slides_dir / "slide_01.html").write_text(
        fixture.read_text(encoding="utf-8"), encoding="utf-8"
    )
    output = tmp_path / f"{slides_dir.name}.pptx"

    await HTMLExporter().export(slides_dir, output, "16:9")
    return Presentation(output)
