from __future__ import annotations

from pathlib import Path
from zipfile import ZipFile

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from localdeck.templates.deck_backend import FrameMap
from localdeck.templates.pptx_backend import PptxTemplateBackend
from tests.fixtures.build_template_deck import build_template_deck


def test_clones_and_edits_source_frame_without_losing_template_parts(
    tmp_path: Path,
) -> None:
    source, replacement_image = build_template_deck(tmp_path)
    source_deck = Presentation(source)
    source_slide = source_deck.slides[1]
    source_title = _shape_by_name(source_slide.shapes, "editable-title")
    source_body = _shape_by_name(source_slide.shapes, "editable-body")
    source_picture = _shape_by_name(source_slide.shapes, "editable-picture")
    picture_box = (
        source_picture.left,
        source_picture.top,
        source_picture.width,
        source_picture.height,
    )
    picture_crop = (
        source_picture.crop_left,
        source_picture.crop_right,
        source_picture.crop_top,
        source_picture.crop_bottom,
    )
    source_theme = _theme_bytes(source)

    backend = PptxTemplateBackend()
    working = tmp_path / "working.pptx"
    backend.create_from_map(source, FrameMap(source_slide_numbers=(2,)), working)
    backend.replace_text(0, source_title.shape_id, "New Editable Title")
    backend.replace_text(0, source_body.shape_id, "New editable body")
    backend.replace_image(0, source_picture.shape_id, replacement_image)
    output = backend.save(tmp_path / "result.pptx")

    generated = Presentation(output)
    assert len(generated.slides) == 1
    slide = generated.slides[0]
    assert slide.slide_layout.name == source_slide.slide_layout.name
    assert _shape_by_name(slide.shapes, "editable-title").text == "New Editable Title"
    assert _shape_by_name(slide.shapes, "editable-body").text == "New editable body"

    picture = _shape_by_name(slide.shapes, "editable-picture")
    assert picture.shape_type == MSO_SHAPE_TYPE.PICTURE
    assert (picture.left, picture.top, picture.width, picture.height) == picture_box
    assert (
        picture.crop_left,
        picture.crop_right,
        picture.crop_top,
        picture.crop_bottom,
    ) == picture_crop
    assert _shape_by_name(slide.shapes, "brand-logo")
    assert _shape_by_name(slide.shapes, "brand-footer")
    assert _shape_by_name(slide.shapes, "page-number")
    assert _shape_by_name(slide.shapes, "preserved-group").shape_type == (
        MSO_SHAPE_TYPE.GROUP
    )
    assert _theme_bytes(output) == source_theme


def _shape_by_name(shapes: object, name: str) -> object:
    return next(shape for shape in shapes if shape.name == name)  # type: ignore[attr-defined]


def _theme_bytes(path: Path) -> bytes:
    with ZipFile(path) as archive:
        return archive.read("ppt/theme/theme1.xml")
