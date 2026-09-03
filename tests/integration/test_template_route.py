from __future__ import annotations

from pathlib import Path
from zipfile import ZipFile

from pptx import Presentation

from localdeck.generation.template_engine.generator import TemplateRouteGenerator
from localdeck.generation.template_engine.matcher import (
    FrameMatchDecision,
    FrameMatchMap,
    ReuseMode,
)
from localdeck.planning.copywriter import SharedCopy
from localdeck.planning.models import ContentBlock, SlidePlan, SlideSpec
from localdeck.templates.inspector import TemplateInspector
from localdeck.templates.models import NarrativeRole
from tests.fixtures.build_template_deck import build_template_deck


def test_generates_editable_source_and_derived_slides_from_template(
    tmp_path: Path,
) -> None:
    source, replacement_image = build_template_deck(tmp_path)
    inspection = TemplateInspector().inspect(source)
    plan = SlidePlan(
        max_slides=30,
        slides=(
            SlideSpec(
                index=1,
                slide_id="cover",
                role=NarrativeRole.COVER,
                title="携手同济大学,共建数智化新生态",
                core_message="华为与同济大学合作交流材料",
                visual_intent="template cover",
            ),
            SlideSpec(
                index=2,
                slide_id="company",
                role=NarrativeRole.CONTENT,
                chapter_index=1,
                section_index=1,
                title="华为经营情况",
                core_message="持续投入根技术,服务产业数智化转型。",
                content_blocks=(
                    ContentBlock(
                        kind="bullets",
                        items=("稳健经营", "持续创新", "开放合作"),
                    ),
                ),
                asset_ids=("hero",),
                visual_intent="title body image",
                source_footer="来源: 华为公司公开资料",
            ),
            SlideSpec(
                index=3,
                slide_id="outlook",
                role=NarrativeRole.SOLUTION,
                chapter_index=4,
                section_index=2,
                title="未来合作建议",
                core_message="围绕科研、人才与产业协同形成长期合作机制。",
                content_blocks=(
                    ContentBlock(
                        kind="bullets",
                        items=("联合创新", "人才培养", "成果转化"),
                    ),
                ),
                asset_ids=("hero",),
                visual_intent="derived template layout",
                source_footer="来源: 双方合作建议",
            ),
        ),
    )
    frame_map = FrameMatchMap(
        decisions=(
            _source_decision("cover", 1, "cover"),
            _source_decision("company", 2, "content"),
            FrameMatchDecision(
                slide_id="outlook",
                reuse_mode=ReuseMode.DERIVED_LAYOUT,
                family="derived-layout",
                score=0,
                score_breakdown={},
            ),
        ),
        omitted_source_slides={},
    )
    output = tmp_path / "template-route.pptx"

    result = TemplateRouteGenerator().generate(
        source=source,
        inspection=inspection,
        shared_copy=SharedCopy(plan=plan, evidence={}),
        frame_map=frame_map,
        output=output,
        assets={"hero": replacement_image},
    )

    assert result == output.resolve()
    generated = Presentation(result)
    assert len(generated.slides) == 3
    assert _theme_bytes(result) == _theme_bytes(source)
    assert generated.slides[0].shapes.title.text == plan.slides[0].title
    assert _text(generated.slides[1]).find("华为经营情况") >= 0
    assert _text(generated.slides[2]).find("未来合作建议") >= 0
    for slide in generated.slides:
        names = {shape.name for shape in slide.shapes}
        assert "brand-logo" in names
        assert "brand-footer" in names
        assert not any(
            shape.is_placeholder and not getattr(shape, "text", "").strip()
            for shape in slide.shapes
        )
    all_text = "\n".join(_text(slide) for slide in generated.slides)
    assert "Template Cover" not in all_text
    assert "Source Title" not in all_text
    assert "Source body" not in all_text
    assert "来源: 华为公司公开资料" in all_text
    assert "来源: 双方合作建议" in all_text


def _source_decision(
    slide_id: str, source_slide_number: int, family: str
) -> FrameMatchDecision:
    return FrameMatchDecision(
        slide_id=slide_id,
        reuse_mode=ReuseMode.SOURCE_FRAME,
        source_slide_number=source_slide_number,
        source_layout_id=f"slide-{source_slide_number:02d}-{family}",
        family=family,
        score=100,
        score_breakdown={"fixture": 100},
    )


def _text(slide: object) -> str:
    return "\n".join(
        getattr(shape, "text", "") for shape in slide.shapes  # type: ignore[attr-defined]
    )


def _theme_bytes(path: Path) -> bytes:
    with ZipFile(path) as archive:
        return archive.read("ppt/theme/theme1.xml")
