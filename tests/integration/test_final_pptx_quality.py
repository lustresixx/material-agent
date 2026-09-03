from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from localdeck.planning.models import SlidePlan, SlideSpec
from localdeck.quality.deck import FinalDeckQualityGate
from localdeck.templates.models import NarrativeRole


def test_rejects_structural_visual_and_brand_failures(tmp_path: Path) -> None:
    path = tmp_path / "bad.pptx"
    presentation = Presentation()
    for index in range(4):
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = f"Slide {index + 1}"
        tiny = slide.shapes.add_textbox(
            Inches(1), Inches(2), Inches(0.5), Inches(0.2)
        )
        tiny.text_frame.text = "This text is much too long for the tiny text box"
        run = tiny.text_frame.paragraphs[0].runs[0]
        run.font.size = Pt(24)
        run.font.color.rgb = RGBColor(255, 255, 255)
    presentation.save(path)
    plan = SlidePlan(
        max_slides=30,
        slides=tuple(
            SlideSpec(
                index=index,
                slide_id=f"slide-{index}",
                role=NarrativeRole.CONTENT,
                chapter_index=1,
                section_index=index,
                title=f"Slide {index}",
                core_message="Core",
                visual_intent="same layout",
                source_footer="Source: official publication",
            )
            for index in range(1, 5)
        ),
    )

    report = FinalDeckQualityGate().inspect(
        path,
        plan=plan,
        required_brand_names=("brand-logo",),
    )

    assert not report.passed
    assert {
        "low-contrast-text",
        "possible-clipped-text",
        "empty-placeholder",
        "missing-brand-furniture",
        "missing-source-footer",
        "repeated-layout-silhouette",
    }.issubset(report.codes)


def test_rejects_more_than_thirty_final_slides(tmp_path: Path) -> None:
    path = tmp_path / "too-many.pptx"
    presentation = Presentation()
    for index in range(31):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        slide.shapes.add_textbox(
            Inches(1), Inches(1), Inches(6), Inches(1)
        ).text_frame.text = f"Slide {index + 1}"
    presentation.save(path)

    report = FinalDeckQualityGate().inspect(path)

    assert "slide-limit" in report.codes


def test_final_render_is_part_of_publication_gate(tmp_path: Path) -> None:
    path = tmp_path / "rendered.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_textbox(
        Inches(1), Inches(1), Inches(6), Inches(1)
    ).text_frame.text = "Visible OOXML text"
    presentation.save(path)

    report = FinalDeckQualityGate().inspect(
        path,
        preview_renderer=_BlankRenderer(),
        previews_dir=tmp_path / "previews",
    )

    assert "blank-final-render" in report.codes


class _BlankRenderer:
    def render(self, source: Path, output_dir: Path) -> list[Path]:
        del source
        output_dir.mkdir(parents=True, exist_ok=True)
        image_path = output_dir / "slide_01.png"
        Image.new("RGB", (320, 180), "white").save(image_path)
        return [image_path]
