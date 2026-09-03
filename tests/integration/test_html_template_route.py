from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation

from localdeck.generation.html_agent.generator import HtmlRouteGenerator
from localdeck.llm.protocol import AssistantResponse
from localdeck.planning.copywriter import SharedCopy
from localdeck.planning.models import ContentBlock, SlidePlan, SlideSpec
from localdeck.templates.models import NarrativeRole, ThemeProfile
from tests.fixtures.scripted_run import ScriptedLLM


async def test_generates_batched_template_constrained_html_and_pptx(
    tmp_path: Path,
) -> None:
    plan = _plan(5)
    llm = ScriptedLLM(
        [
            AssistantResponse(content=_batch_response(plan.slides[:3])),
            AssistantResponse(content=_batch_response(plan.slides[3:])),
        ]
    )
    output = tmp_path / "html-route.pptx"

    result = await HtmlRouteGenerator(llm, batch_size=3).generate(
        shared_copy=SharedCopy(plan=plan, evidence={}),
        theme=_theme(),
        workspace=tmp_path / "html-work",
        output=output,
        aspect_ratio="16:9",
    )

    assert result.pptx == output.resolve()
    assert result.model_calls == 2
    assert result.repairs == 0
    assert result.fallback_slides == ()
    assert len(Presentation(result.pptx).slides) == 5
    assert len(list(result.slides_dir.glob("slide_*.html"))) == 5
    theme_css = (result.slides_dir / "theme.css").read_text(encoding="utf-8")
    assert "Microsoft YaHei" in theme_css
    assert "#E60012" in theme_css


async def test_repairs_only_the_failed_slide(tmp_path: Path) -> None:
    plan = _plan(2)
    bad = _slide_html(plan.slides[1]).replace(
        "</body>", '<script src="https://example.com/x.js"></script></body>'
    )
    first = json.dumps(
        {
            "slides": [
                {
                    "slide_id": plan.slides[0].slide_id,
                    "html": _slide_html(plan.slides[0]),
                },
                {"slide_id": plan.slides[1].slide_id, "html": bad},
            ]
        },
        ensure_ascii=False,
    )
    llm = ScriptedLLM(
        [
            AssistantResponse(content=first),
            AssistantResponse(content=_batch_response(plan.slides[1:])),
        ]
    )

    result = await HtmlRouteGenerator(llm, batch_size=2, max_repairs=1).generate(
        shared_copy=SharedCopy(plan=plan, evidence={}),
        theme=_theme(),
        workspace=tmp_path / "repair-work",
        output=tmp_path / "repaired.pptx",
        aspect_ratio="16:9",
    )

    assert result.repairs == 1
    assert result.model_calls == 2
    repair_prompt = str(llm.requests[1][-1]["content"])
    assert plan.slides[1].slide_id in repair_prompt
    assert plan.slides[0].slide_id not in repair_prompt


async def test_uses_safe_fallback_after_repair_budget_is_exhausted(
    tmp_path: Path,
) -> None:
    plan = _plan(2)
    invalid_response = AssistantResponse(content="not-json")
    llm = ScriptedLLM([invalid_response, invalid_response])

    result = await HtmlRouteGenerator(llm, batch_size=2, max_repairs=1).generate(
        shared_copy=SharedCopy(plan=plan, evidence={}),
        theme=_theme(),
        workspace=tmp_path / "fallback-work",
        output=tmp_path / "fallback.pptx",
        aspect_ratio="16:9",
    )

    assert result.repairs == 1
    assert result.fallback_slides == ("slide-01", "slide-02")
    assert len(Presentation(result.pptx).slides) == 2


def _plan(count: int) -> SlidePlan:
    return SlidePlan(
        max_slides=30,
        slides=tuple(
            SlideSpec(
                index=index,
                slide_id=f"slide-{index:02d}",
                role=NarrativeRole.CONTENT,
                chapter_index=1,
                section_index=index,
                title=f"第 {index} 页",
                core_message=f"核心信息 {index}",
                content_blocks=(
                    ContentBlock(kind="bullets", items=("要点一", "要点二")),
                ),
                visual_intent="template constrained html",
                source_footer="来源: 官方公开资料",
            )
            for index in range(1, count + 1)
        ),
    )


def _theme() -> ThemeProfile:
    return ThemeProfile(
        page_width=13.333,
        page_height=7.5,
        font_families=("Microsoft YaHei",),
        palette=("#E60012", "#111111", "#FFFFFF"),
        spacing=(0.08, 0.16, 0.24, 0.4, 0.6),
    )


def _batch_response(slides: tuple[SlideSpec, ...]) -> str:
    return json.dumps(
        {
            "slides": [
                {"slide_id": slide.slide_id, "html": _slide_html(slide)}
                for slide in slides
            ]
        },
        ensure_ascii=False,
    )


def _slide_html(slide: SlideSpec) -> str:
    return f"""<!doctype html>
<html lang="zh-CN"><head><meta charset="utf-8">
<link rel="stylesheet" href="global.css"><link rel="stylesheet" href="theme.css">
</head><body><main><h1>{slide.title}</h1><p>{slide.core_message}</p></main>
<footer>{slide.source_footer}</footer></body></html>"""
