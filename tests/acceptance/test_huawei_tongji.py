from __future__ import annotations

import json
import os
from pathlib import Path

import pytest
from pptx import Presentation

from localdeck.config import Settings
from localdeck.models import GenerationRoute, TemplateGenerationRequest
from localdeck.pipeline import TemplateDeckPipeline

pytestmark = pytest.mark.live

OUTLINE = Path(__file__).parents[2] / "examples" / "huawei-tongji-outline.json"


@pytest.mark.skipif(
    os.getenv("LOCALDECK_RUN_LIVE_TESTS") != "1"
    or not os.getenv("LOCALDECK_ACCEPTANCE_TEMPLATE"),
    reason=(
        "Set LOCALDECK_RUN_LIVE_TESTS=1 and LOCALDECK_ACCEPTANCE_TEMPLATE "
        "to run the real Huawei/Tongji acceptance scenario"
    ),
)
async def test_huawei_tongji_dual_route_acceptance(tmp_path: Path) -> None:
    template = Path(os.environ["LOCALDECK_ACCEPTANCE_TEMPLATE"]).resolve()
    request = TemplateGenerationRequest(
        outline=OUTLINE,
        template=str(template),
        routes=(GenerationRoute.TEMPLATE, GenerationRoute.HTML),
        max_slides=30,
        output_dir=tmp_path / "published",
    )

    result = await TemplateDeckPipeline(Settings.from_env()).generate(request)

    manifest = json.loads(result.manifest.read_text(encoding="utf-8"))
    plan = json.loads(result.plan.read_text(encoding="utf-8"))
    planned_sections = {
        (slide["chapter_index"], slide["section_index"])
        for slide in plan["slides"]
        if slide["chapter_index"] is not None and slide["section_index"] is not None
    }
    expected_sections = {
        (chapter_index, section_index)
        for chapter_index, chapter in enumerate(
            json.loads(OUTLINE.read_text(encoding="utf-8"))["chapters"], start=1
        )
        for section_index, _ in enumerate(chapter["sections"], start=1)
    }

    assert "docker" not in json.dumps(manifest).casefold()
    assert template.is_file() and OUTLINE.is_file()
    assert planned_sections == expected_sections
    assert len(plan["slides"]) <= 30
    assert expected_sections.issubset(planned_sections)
    assert plan["slides"][0]["role"] == "cover"
    assert list((result.workspace / "research").rglob("sources.json"))
    assert any(slide["source_footer"] for slide in plan["slides"])
    assert manifest["routes"]["template"]["plan_digest"] == manifest["plan_digest"]
    assert manifest["routes"]["html"]["plan_digest"] == manifest["plan_digest"]
    assert (result.workspace / "planning" / "frame-map.json").is_file()
    assert (result.workspace / "routes" / "html" / "slides" / "theme.css").is_file()
    assert result.template_output is not None and result.template_output.is_file()
    assert result.html_output is not None and result.html_output.is_file()
    assert all(
        len(Presentation(str(path)).slides) == len(plan["slides"])
        for path in (result.template_output, result.html_output)
    )
    assert manifest["routes"]["template"]["status"] == "completed"
    assert all(
        "brand-logo" in {shape.name for shape in slide.shapes}
        for slide in Presentation(str(result.template_output)).slides
    )
    assert "low-contrast-text" not in manifest["routes"]["html"]["quality_issues"]
    assert result.comparison is not None and result.comparison.is_file()
    comparison = result.comparison.read_text(encoding="utf-8")
    assert "model calls:" in comparison and "duration:" in comparison
    assert not list(request.output_dir.glob(".*.tmp"))
