"""Build a deterministic editable PPTX fixture for template backend tests."""

from __future__ import annotations

from pathlib import Path

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt


def build_template_deck(directory: Path) -> tuple[Path, Path]:
    """Create a two-slide source deck and a replacement image."""
    directory.mkdir(parents=True, exist_ok=True)
    logo_path = directory / "logo.png"
    source_image_path = directory / "source-image.png"
    replacement_image_path = directory / "replacement-image.png"
    _write_image(logo_path, (20, 90, 220), "LOGO", (320, 100))
    _write_image(source_image_path, (20, 150, 110), "SOURCE", (800, 500))
    _write_image(replacement_image_path, (220, 80, 20), "NEW", (600, 900))

    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)

    cover = presentation.slides.add_slide(presentation.slide_layouts[0])
    cover.shapes.title.text = "Template Cover"
    _add_brand_furniture(cover, logo_path, "01")

    content = presentation.slides.add_slide(presentation.slide_layouts[1])
    title = content.shapes.title
    title.name = "editable-title"
    title.text = "Source Title"
    body = content.placeholders[1]
    body.name = "editable-body"
    body.text = "Source body"

    picture = content.shapes.add_picture(
        str(source_image_path), Inches(8.2), Inches(1.6), Inches(4.2), Inches(4.5)
    )
    picture.name = "editable-picture"
    picture.crop_left = 0.08
    picture.crop_right = 0.12
    picture.crop_top = 0.05
    picture.crop_bottom = 0.09

    group = content.shapes.add_group_shape()
    group.name = "preserved-group"
    group.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(0.8),
        Inches(5.8),
        Inches(0.35),
        Inches(0.35),
    )
    group.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1.2),
        Inches(5.9),
        Inches(1.6),
        Inches(0.15),
    )
    _add_brand_furniture(content, logo_path, "02")

    source = directory / "template-source.pptx"
    presentation.save(source)
    return source, replacement_image_path


def _add_brand_furniture(slide: object, logo_path: Path, page_number: str) -> None:
    shapes = slide.shapes  # type: ignore[attr-defined]
    logo = shapes.add_picture(
        str(logo_path), Inches(11.7), Inches(0.25), Inches(1.0), Inches(0.3)
    )
    logo.name = "brand-logo"
    footer = shapes.add_textbox(
        Inches(0.6), Inches(7.05), Inches(6.0), Inches(0.25)
    )
    footer.name = "brand-footer"
    footer.text_frame.text = "Huawei enterprise cooperation material"
    footer.text_frame.paragraphs[0].font.size = Pt(8)
    page = shapes.add_textbox(
        Inches(12.3), Inches(7.02), Inches(0.4), Inches(0.25)
    )
    page.name = "page-number"
    page.text_frame.text = page_number
    page.text_frame.paragraphs[0].font.size = Pt(8)


def _write_image(
    path: Path,
    color: tuple[int, int, int],
    label: str,
    size: tuple[int, int],
) -> None:
    image = Image.new("RGB", size, color)
    draw = ImageDraw.Draw(image)
    draw.text((20, 20), label, fill=(255, 255, 255))
    image.save(path)
