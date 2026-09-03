"""Static final-PPTX checks independent of the HTML source path."""

from __future__ import annotations

from collections.abc import Iterable
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor

from localdeck.planning.models import SlidePlan
from localdeck.quality.deck import QualityIssue

_EMU_PER_INCH = 914400


def inspect_pptx(
    path: Path,
    *,
    plan: SlidePlan | None = None,
    required_brand_names: tuple[str, ...] = (),
) -> tuple[QualityIssue, ...]:
    """Inspect editable OOXML content and conservative geometry heuristics."""
    presentation = Presentation(str(path.expanduser().resolve()))
    issues: list[QualityIssue] = []
    if len(presentation.slides) > 30:
        issues.append(
            QualityIssue(
                code="slide-limit",
                message=(
                    f"Final deck has {len(presentation.slides)} slides; "
                    "maximum is 30"
                ),
            )
        )
    if plan is not None and len(presentation.slides) != len(plan.slides):
        issues.append(
            QualityIssue(
                code="slide-count",
                message="Final deck slide count differs from the accepted plan",
            )
        )

    for slide_index, slide in enumerate(presentation.slides, start=1):
        names = {shape.name for shape in slide.shapes}
        for required_name in required_brand_names:
            if required_name not in names:
                issues.append(
                    QualityIssue(
                        code="missing-brand-furniture",
                        message=f"Required brand shape is missing: {required_name}",
                        slide_index=slide_index,
                    )
                )
        for shape in _walk_shapes(slide.shapes):
            text = getattr(shape, "text", "").strip()
            if getattr(shape, "is_placeholder", False) and not text:
                issues.append(
                    QualityIssue(
                        code="empty-placeholder",
                        message="Inherited empty placeholder remains in final PPTX",
                        slide_index=slide_index,
                    )
                )
            if not text or not getattr(shape, "has_text_frame", False):
                continue
            if _has_white_text(shape):
                issues.append(
                    QualityIssue(
                        code="low-contrast-text",
                        message=(
                            "White text may be invisible on a lost/default background"
                        ),
                        slide_index=slide_index,
                    )
                )
            if _likely_clipped(shape, text):
                issues.append(
                    QualityIssue(
                        code="possible-clipped-text",
                        message=(
                            "Text density exceeds the shape's conservative capacity"
                        ),
                        slide_index=slide_index,
                    )
                )
        if plan is not None and slide_index <= len(plan.slides):
            expected_footer = plan.slides[slide_index - 1].source_footer
            visible = "\n".join(
                getattr(shape, "text", "") for shape in _walk_shapes(slide.shapes)
            )
            if expected_footer and expected_footer not in visible:
                issues.append(
                    QualityIssue(
                        code="missing-source-footer",
                        message="Planned source footer is absent from final PPTX",
                        slide_index=slide_index,
                    )
                )
    return tuple(issues)


def _walk_shapes(shapes: Iterable[Any]) -> list[Any]:
    result: list[Any] = []
    for shape in shapes:
        result.append(shape)
        if hasattr(shape, "shapes"):
            result.extend(_walk_shapes(shape.shapes))
    return result


def _has_white_text(shape: Any) -> bool:
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            try:
                if run.font.color.rgb == RGBColor(255, 255, 255):
                    return True
            except (AttributeError, TypeError):
                continue
    return False


def _likely_clipped(shape: Any, text: str) -> bool:
    width_inches = max(float(shape.width) / _EMU_PER_INCH, 0.1)
    height_inches = max(float(shape.height) / _EMU_PER_INCH, 0.1)
    sizes = [
        float(run.font.size.pt)
        for paragraph in shape.text_frame.paragraphs
        for run in paragraph.runs
        if run.font.size is not None
    ]
    font_points = max(sizes, default=18.0)
    line_height_inches = font_points * 1.25 / 72
    chars_per_line = max(int(width_inches * 72 / (font_points * 0.55)), 1)
    available_lines = max(int(height_inches / line_height_inches), 1)
    return len(text) > chars_per_line * available_lines
