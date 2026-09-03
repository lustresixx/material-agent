"""Aggregated deck-level quality checks and publication receipt."""

from __future__ import annotations

from collections import Counter
from pathlib import Path
from typing import TYPE_CHECKING, Any, cast

from PIL import Image, ImageStat
from pptx import Presentation
from pydantic import BaseModel, ConfigDict

from localdeck.planning.models import SlidePlan

if TYPE_CHECKING:
    from localdeck.rendering.pptx_preview import PPTXPreviewRenderer


class QualityIssue(BaseModel):
    """One actionable final-deck quality failure."""

    model_config = ConfigDict(frozen=True)

    code: str
    message: str
    slide_index: int | None = None
    severity: str = "error"


class QualityReport(BaseModel):
    """Combined quality result suitable for manifests and comparison reports."""

    model_config = ConfigDict(frozen=True)

    passed: bool
    issues: tuple[QualityIssue, ...] = ()

    @property
    def codes(self) -> set[str]:
        """Return distinct issue codes for concise assertions and diagnostics."""
        return {issue.code for issue in self.issues}


class FinalDeckQualityGate:
    """Combine final-PPTX checks with whole-deck composition checks."""

    def inspect(
        self,
        path: Path,
        *,
        plan: SlidePlan | None = None,
        required_brand_names: tuple[str, ...] = (),
        preview_renderer: PPTXPreviewRenderer | None = None,
        previews_dir: Path | None = None,
    ) -> QualityReport:
        """Return a publication receipt; callers decide whether to publish."""
        from localdeck.quality.pptx import inspect_pptx

        issues = list(
            inspect_pptx(
                path,
                plan=plan,
                required_brand_names=required_brand_names,
            )
        )
        presentation = Presentation(str(path.expanduser().resolve()))
        if preview_renderer is not None:
            if previews_dir is None:
                raise ValueError("previews_dir is required with preview_renderer")
            previews = preview_renderer.render(path, previews_dir)
            if len(previews) != len(presentation.slides):
                issues.append(
                    QualityIssue(
                        code="final-render-count",
                        message=(
                            "Final PPTX renderer returned an incomplete preview set"
                        ),
                    )
                )
            for slide_index, preview in enumerate(previews, start=1):
                if _is_blank_render(preview):
                    issues.append(
                        QualityIssue(
                            code="blank-final-render",
                            message="Final PPTX render is visually blank",
                            slide_index=slide_index,
                        )
                    )
        signatures = [_silhouette(slide) for slide in presentation.slides]
        if len(signatures) >= 4:
            _, count = Counter(signatures).most_common(1)[0]
            if count / len(signatures) >= 0.8:
                issues.append(
                    QualityIssue(
                        code="repeated-layout-silhouette",
                        message="At least 80% of slides repeat one layout silhouette",
                        severity="warning",
                    )
                )
        deduplicated = tuple(
            {
                (issue.code, issue.slide_index, issue.message): issue
                for issue in issues
            }.values()
        )
        return QualityReport(
            passed=not any(issue.severity == "error" for issue in deduplicated),
            issues=deduplicated,
        )


def _silhouette(slide: Any) -> tuple[tuple[int, int, int, int, int], ...]:
    return tuple(
        sorted(
            (
                int(shape.shape_type),
                round(int(shape.left) / 914400 * 10),
                round(int(shape.top) / 914400 * 10),
                round(int(shape.width) / 914400 * 10),
                round(int(shape.height) / 914400 * 10),
            )
            for shape in slide.shapes
        )
    )


def _is_blank_render(path: Path) -> bool:
    with Image.open(path) as source:
        image = source.convert("RGB")
        statistics = ImageStat.Stat(image)
    extrema = cast(tuple[tuple[int, int], ...], image.getextrema())
    channel_ranges = [high - low for low, high in extrema]
    return max(channel_ranges) <= 3 and min(statistics.mean) >= 248
