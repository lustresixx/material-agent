"""Reuse approved template frames and furniture for derived slides."""

from __future__ import annotations

from copy import deepcopy
from typing import Any, cast

from pptx import Presentation
from pptx.presentation import Presentation as PresentationType
from pptx.slide import Slide

from localdeck.planning.models import SlideSpec
from localdeck.templates.models import LayoutFrame, NarrativeRole


def select_derived_frame(
    slide: SlideSpec, frames: tuple[LayoutFrame, ...]
) -> LayoutFrame:
    """Choose a template-owned frame without inventing a new style system."""
    if not frames:
        raise ValueError("template inspection contains no reusable frames")
    same_role = [frame for frame in frames if frame.role == slide.role]
    content = [frame for frame in frames if frame.role == NarrativeRole.CONTENT]
    candidates = same_role or content or list(frames)
    return max(
        candidates,
        key=lambda frame: (
            len(frame.editable_slots),
            frame.capacity.max_characters,
            -frame.source_slide_number,
        ),
    )


def finalize_template_furniture(
    path: Any, source_footers: tuple[str | None, ...]
) -> None:
    """Add template-styled citations, update page numbers, and remove blanks."""
    presentation: PresentationType = Presentation(str(path))
    for page_number, (slide, source_footer) in enumerate(
        zip(presentation.slides, source_footers, strict=True), start=1
    ):
        _set_page_number(slide, page_number)
        if source_footer:
            _clone_source_footer(slide, source_footer)
        _remove_empty_placeholders(slide)
    presentation.save(str(path))


def _set_page_number(slide: Slide, page_number: int) -> None:
    for shape in slide.shapes:
        if shape.name == "page-number" and shape.has_text_frame:
            _replace_text(cast(Any, shape).text_frame, f"{page_number:02d}")


def _clone_source_footer(slide: Slide, text: str) -> None:
    donor = next(
        (
            shape
            for shape in slide.shapes
            if shape.name == "brand-footer" and shape.has_text_frame
        ),
        None,
    )
    if donor is None:
        return
    element = deepcopy(donor.element)
    properties = element.xpath(".//p:cNvPr")[0]
    properties.set("id", str(_next_shape_id(slide)))
    properties.set("name", "source-footer")
    offset = element.xpath(".//a:xfrm/a:off")
    if offset:
        y = max(0, int(offset[0].get("y")) - int(donor.height) - 45720)
        offset[0].set("y", str(y))
    tree = slide.shapes._spTree  # pyright: ignore[reportPrivateUsage]
    tree.insert_element_before(element, "p:extLst")
    added = next(shape for shape in slide.shapes if shape.name == "source-footer")
    _replace_text(cast(Any, added).text_frame, text)


def _replace_text(text_frame: Any, text: str) -> None:
    paragraph = text_frame.paragraphs[0]
    if paragraph.runs:
        paragraph.runs[0].text = text
        for run in paragraph.runs[1:]:
            paragraph._p.remove(run._r)  # pyright: ignore[reportPrivateUsage]
    else:
        paragraph.add_run().text = text
    for extra in text_frame.paragraphs[1:]:
        text_frame._txBody.remove(extra._p)  # pyright: ignore[reportPrivateUsage]


def _next_shape_id(slide: Slide) -> int:
    return max((shape.shape_id for shape in slide.shapes), default=1) + 1


def _remove_empty_placeholders(slide: Slide) -> None:
    tree = slide.shapes._spTree  # pyright: ignore[reportPrivateUsage]
    for shape in list(slide.shapes):
        if shape.is_placeholder and not getattr(shape, "text", "").strip():
            tree.remove(shape.element)
