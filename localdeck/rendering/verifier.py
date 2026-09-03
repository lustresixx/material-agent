"""Structural verification for generated PowerPoint files."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pydantic import BaseModel

from localdeck.planning.models import SlidePlan
from localdeck.quality.deck import FinalDeckQualityGate


class PPTXVerificationError(ValueError):
    """Raised when a generated PPTX cannot satisfy the output contract."""


class PPTXVerification(BaseModel):
    """Small verification receipt persisted by the pipeline."""

    path: Path
    slide_count: int
    shape_count: int
    text_shape_count: int


class PPTXVerifier:
    """Reopen a PPTX and verify the minimum editable-content guarantees."""

    def verify(self, path: Path, *, expected_slides: int) -> PPTXVerification:
        """Return structural counts or raise before the file is published."""

        if not path.is_file() or path.stat().st_size == 0:
            raise PPTXVerificationError(f"PPTX does not exist or is empty: {path}")
        try:
            presentation = Presentation(str(path))
        except Exception as error:
            raise PPTXVerificationError(f"PPTX cannot be opened: {error}") from error

        slide_count = len(presentation.slides)
        if slide_count != expected_slides:
            raise PPTXVerificationError(
                f"Expected {expected_slides} slides; found {slide_count}"
            )

        shapes = [shape for slide in presentation.slides for shape in slide.shapes]
        text_shapes = []
        for shape in shapes:
            text = getattr(shape, "text", "")
            if getattr(shape, "has_text_frame", False) and text.strip():
                text_shapes.append(shape)
        if not shapes or not text_shapes:
            raise PPTXVerificationError("PPTX contains no editable text content")
        return PPTXVerification(
            path=path,
            slide_count=slide_count,
            shape_count=len(shapes),
            text_shape_count=len(text_shapes),
        )

    def verify_final(
        self,
        path: Path,
        *,
        plan: SlidePlan,
        required_brand_names: tuple[str, ...] = (),
    ) -> PPTXVerification:
        """Require both structural validity and final-deck quality gates."""
        verification = self.verify(path, expected_slides=len(plan.slides))
        report = FinalDeckQualityGate().inspect(
            path,
            plan=plan,
            required_brand_names=required_brand_names,
        )
        if not report.passed:
            codes = ", ".join(sorted(report.codes))
            raise PPTXVerificationError(f"Final PPTX quality failed: {codes}")
        return verification
