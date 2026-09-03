"""Editable template backend built on python-pptx and source slide cloning."""

from __future__ import annotations

import re
from pathlib import Path
from typing import Any, cast

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.presentation import Presentation as PresentationType
from pptx.shapes.base import BaseShape
from pptx.shapes.group import GroupShape
from pptx.shapes.picture import Picture
from pptx.slide import Slide

from localdeck.templates.deck_backend import FrameMap
from localdeck.templates.models import TemplateManifest
from localdeck.templates.relationships import clone_slide


class PptxTemplateBackend:
    """Clone and edit the common PowerPoint shape subset without Office."""

    def __init__(self) -> None:
        self._presentation: PresentationType | None = None

    def inspect(self, source: Path) -> TemplateManifest:
        """Return source identity and slide count without modifying the deck."""
        path = source.expanduser().resolve()
        presentation = Presentation(str(path))
        template_id = re.sub(r"[^a-z0-9]+", "-", path.stem.casefold()).strip("-")
        return TemplateManifest(
            template_id=template_id or "template",
            name=path.stem,
            source_file=path.name,
            slide_count=len(presentation.slides),
        )

    def create_from_map(
        self, source: Path, frame_map: FrameMap, output: Path
    ) -> Path:
        """Clone requested source frames and discard sample narrative slides."""
        source_path = source.expanduser().resolve()
        presentation = Presentation(str(source_path))
        source_count = len(presentation.slides)
        for slide_number in frame_map.source_slide_numbers:
            if slide_number > source_count:
                raise IndexError(f"source slide number out of range: {slide_number}")
            clone_slide(presentation, presentation.slides[slide_number - 1])

        for _ in range(source_count):
            _remove_slide(presentation, 0)

        self._presentation = presentation
        return self.save(output)

    def replace_text(self, slide_index: int, shape_id: int, text: str) -> None:
        """Replace visible text while retaining the first run's formatting."""
        shape = self._find_shape(slide_index, shape_id)
        if not shape.has_text_frame:
            raise TypeError(f"shape {shape_id} does not contain editable text")
        text_frame = cast(Any, shape).text_frame
        first_paragraph = text_frame.paragraphs[0]
        if first_paragraph.runs:
            first_run = first_paragraph.runs[0]
            first_run.text = text
            for run in first_paragraph.runs[1:]:
                first_paragraph._p.remove(run._r)  # pyright: ignore[reportPrivateUsage]
        else:
            first_paragraph.add_run().text = text
        for paragraph in text_frame.paragraphs[1:]:
            text_frame._txBody.remove(  # pyright: ignore[reportPrivateUsage]
                paragraph._p  # pyright: ignore[reportPrivateUsage]
            )

    def replace_image(self, slide_index: int, shape_id: int, image: Path) -> None:
        """Replace a picture in place and retain crop, name, ID, and z-order."""
        shape = self._find_shape(slide_index, shape_id)
        if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
            raise TypeError(f"shape {shape_id} is not a picture")
        old_picture = cast(Picture, shape)
        slide = self._require_presentation().slides[slide_index]
        new_picture = slide.shapes.add_picture(
            str(image.expanduser().resolve()),
            old_picture.left,
            old_picture.top,
            old_picture.width,
            old_picture.height,
        )
        new_picture.crop_left = old_picture.crop_left
        new_picture.crop_right = old_picture.crop_right
        new_picture.crop_top = old_picture.crop_top
        new_picture.crop_bottom = old_picture.crop_bottom

        old_element = old_picture.element
        parent = old_element.getparent()
        position = parent.index(old_element)
        parent.remove(old_element)
        new_parent = new_picture.element.getparent()
        new_parent.remove(new_picture.element)
        parent.insert(position, new_picture.element)

        properties = new_picture.element.xpath(".//p:cNvPr")[0]
        properties.set("id", str(shape_id))
        properties.set("name", old_picture.name)

    def save(self, output: Path) -> Path:
        """Atomically write the current editable presentation to disk."""
        presentation = self._require_presentation()
        destination = output.expanduser().resolve()
        destination.parent.mkdir(parents=True, exist_ok=True)
        presentation.save(str(destination))
        return destination

    def _find_shape(self, slide_index: int, shape_id: int) -> BaseShape:
        presentation = self._require_presentation()
        if not 0 <= slide_index < len(presentation.slides):
            raise IndexError(f"slide index out of range: {slide_index}")
        shape = _find_shape_by_id(presentation.slides[slide_index], shape_id)
        if shape is None:
            raise KeyError(f"shape ID not found: {shape_id}")
        return shape

    def _require_presentation(self) -> PresentationType:
        if self._presentation is None:
            raise RuntimeError("template backend has no working presentation")
        return self._presentation


def _find_shape_by_id(slide: Slide, shape_id: int) -> BaseShape | None:
    return _find_in_collection(slide.shapes, shape_id)


def _find_in_collection(shapes: Any, shape_id: int) -> BaseShape | None:
    for shape in shapes:
        if shape.shape_id == shape_id:
            return cast(BaseShape, shape)
        if isinstance(shape, GroupShape):
            nested = _find_in_collection(shape.shapes, shape_id)
            if nested is not None:
                return nested
    return None


def _remove_slide(presentation: PresentationType, slide_index: int) -> None:
    slides: Any = presentation.slides
    slide_id = slides._sldIdLst[slide_index]
    presentation.part.drop_rel(slide_id.rId)
    slides._sldIdLst.remove(slide_id)
