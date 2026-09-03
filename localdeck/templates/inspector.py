"""Deterministic full-deck inventory for reusable PowerPoint templates."""

from __future__ import annotations

import hashlib
import re
from collections import Counter
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER

from localdeck.templates.classifier import classify_slide
from localdeck.templates.models import (
    CapacityProfile,
    ComponentSpec,
    EditableSlot,
    EditPolicy,
    LayoutFrame,
    SlotType,
    TemplateInspection,
    TemplateManifest,
)
from localdeck.templates.theme import extract_theme

_EMU_PER_INCH = 914400


class TemplateInspector:
    """Inspect every source slide and retain stable source-shape identities."""

    def inspect(self, source: Path) -> TemplateInspection:
        """Return a complete inventory without changing the source deck."""
        path = source.expanduser().resolve()
        presentation = Presentation(str(path))
        recurring_names = _recurring_shape_names(presentation)
        layouts: list[LayoutFrame] = []
        components: list[ComponentSpec] = []
        for slide_number, slide in enumerate(presentation.slides, start=1):
            slots: list[EditableSlot] = []
            preserve_ids: list[int] = []
            source_ids: list[int] = []
            text_parts: list[str] = []
            image_count = 0
            for shape in _walk_shapes(slide.shapes):
                source_ids.append(shape.shape_id)
                text = getattr(shape, "text", "").strip()
                if text:
                    text_parts.append(text)
                slot_type, edit_policy = _classify_shape(shape, recurring_names)
                if edit_policy == EditPolicy.PRESERVE_ONLY:
                    preserve_ids.append(shape.shape_id)
                    continue
                if slot_type == SlotType.IMAGE:
                    image_count += 1
                capacity = _capacity(shape, slot_type)
                geometry = _geometry(shape)
                slot = EditableSlot(
                    source_shape_id=shape.shape_id,
                    name=shape.name,
                    slot_type=slot_type,
                    edit_policy=edit_policy,
                    capacity=capacity,
                    **geometry,
                )
                slots.append(slot)
                components.append(
                    ComponentSpec(
                        component_id=f"s{slide_number}-shape-{shape.shape_id}",
                        source_slide_number=slide_number,
                        source_shape_ids=(shape.shape_id,),
                        slot_type=slot_type,
                        edit_policy=edit_policy,
                        capacity=capacity,
                        **geometry,
                    )
                )
            body_count = sum(slot.slot_type == SlotType.BODY for slot in slots)
            role, family, confidence = classify_slide(
                slide_number=slide_number,
                text="\n".join(text_parts),
                image_count=image_count,
                body_slot_count=body_count,
            )
            layouts.append(
                LayoutFrame(
                    layout_id=f"slide-{slide_number:02d}-{role.value}",
                    source_slide_number=slide_number,
                    role=role,
                    family=family,
                    capacity=CapacityProfile(
                        max_characters=sum(
                            slot.capacity.max_characters for slot in slots
                        ),
                        max_items=sum(slot.capacity.max_items for slot in slots),
                        max_images=sum(slot.capacity.max_images for slot in slots),
                    ),
                    editable_slots=tuple(slots),
                    preserve_shape_ids=tuple(dict.fromkeys(preserve_ids)),
                    source_shape_ids=tuple(dict.fromkeys(source_ids)),
                    classification_confidence=confidence,
                )
            )
        digest = hashlib.sha256(path.read_bytes()).hexdigest()[:10]
        slug = re.sub(r"[^a-z0-9]+", "-", path.stem.casefold()).strip("-")
        manifest = TemplateManifest(
            template_id=f"{slug or 'template'}-{digest}",
            name=path.stem,
            source_file="source.pptx",
            slide_count=len(presentation.slides),
        )
        return TemplateInspection(
            manifest=manifest,
            theme=extract_theme(presentation),
            layouts=tuple(layouts),
            components=tuple(components),
        )


def _recurring_shape_names(presentation: Any) -> set[str]:
    counts: Counter[str] = Counter()
    for slide in presentation.slides:
        counts.update({shape.name.casefold() for shape in slide.shapes})
    threshold = max(2, (len(presentation.slides) + 1) // 2)
    return {name for name, count in counts.items() if count >= threshold}


def _walk_shapes(shapes: Any) -> list[Any]:
    result: list[Any] = []
    for shape in shapes:
        result.append(shape)
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            result.extend(_walk_shapes(shape.shapes))
    return result


def _classify_shape(
    shape: Any, recurring_names: set[str]
) -> tuple[SlotType, EditPolicy]:
    name = shape.name.casefold()
    if name in recurring_names or shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        return SlotType.OTHER, EditPolicy.PRESERVE_ONLY
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        return SlotType.IMAGE, EditPolicy.REPLACE_IMAGE
    if getattr(shape, "has_text_frame", False):
        if getattr(shape, "is_placeholder", False):
            placeholder_type = shape.placeholder_format.type
            if placeholder_type in {PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE}:
                return SlotType.TITLE, EditPolicy.REPLACE_TEXT
        if "source" in name or "来源" in getattr(shape, "text", ""):
            return SlotType.SOURCE, EditPolicy.REPLACE_TEXT
        return SlotType.BODY, EditPolicy.REPLACE_TEXT
    return SlotType.OTHER, EditPolicy.PRESERVE_ONLY


def _geometry(shape: Any) -> dict[str, float]:
    return {
        "x": round(shape.left / _EMU_PER_INCH, 4),
        "y": round(shape.top / _EMU_PER_INCH, 4),
        "width": round(shape.width / _EMU_PER_INCH, 4),
        "height": round(shape.height / _EMU_PER_INCH, 4),
    }


def _capacity(shape: Any, slot_type: SlotType) -> CapacityProfile:
    if slot_type == SlotType.IMAGE:
        return CapacityProfile(max_images=1)
    area = max((shape.width / _EMU_PER_INCH) * (shape.height / _EMU_PER_INCH), 0)
    return CapacityProfile(
        max_characters=max(12, round(area * 18)),
        max_items=max(1, round(shape.height / _EMU_PER_INCH / 0.55)),
    )
