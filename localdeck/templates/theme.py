"""Extract constrained design tokens from editable PowerPoint content."""

from __future__ import annotations

from typing import Any

from pptx.presentation import Presentation as PresentationType

from localdeck.templates.models import ThemeProfile

_EMU_PER_INCH = 914400


def extract_theme(presentation: PresentationType) -> ThemeProfile:
    """Collect page geometry, explicit fonts, and explicit RGB colors."""
    slide_width = presentation.slide_width
    slide_height = presentation.slide_height
    if slide_width is None or slide_height is None:
        raise ValueError("presentation page geometry is unavailable")
    fonts: set[str] = set()
    colors: set[str] = set()
    for slide in presentation.slides:
        for shape in _walk_shapes(slide.shapes):
            if getattr(shape, "has_text_frame", False):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.name:
                            fonts.add(run.font.name)
                        color = _rgb(run.font.color)
                        if color:
                            colors.add(color)
            fill_color = _fill_rgb(getattr(shape, "fill", None))
            if fill_color:
                colors.add(fill_color)
            line_color = _rgb(getattr(getattr(shape, "line", None), "color", None))
            if line_color:
                colors.add(line_color)
    return ThemeProfile(
        page_width=round(slide_width / _EMU_PER_INCH, 3),
        page_height=round(slide_height / _EMU_PER_INCH, 3),
        font_families=tuple(sorted(fonts)) or ("Microsoft YaHei",),
        palette=tuple(sorted(colors)) or ("#000000", "#FFFFFF"),
        spacing=(0.08, 0.16, 0.24, 0.4, 0.6),
    )


def _walk_shapes(shapes: Any) -> list[Any]:
    result: list[Any] = []
    for shape in shapes:
        result.append(shape)
        if hasattr(shape, "shapes"):
            result.extend(_walk_shapes(shape.shapes))
    return result


def _rgb(color_format: Any) -> str | None:
    if color_format is None:
        return None
    try:
        value = color_format.rgb
    except (AttributeError, TypeError):
        return None
    return f"#{value}" if value is not None else None


def _fill_rgb(fill_format: Any) -> str | None:
    if fill_format is None:
        return None
    try:
        return _rgb(fill_format.fore_color)
    except (AttributeError, TypeError):
        return None
