"""Low-level relationship-aware cloning for python-pptx slide parts."""

from __future__ import annotations

from copy import deepcopy
from typing import Any

from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.presentation import Presentation as PresentationType
from pptx.slide import Slide


def clone_slide(presentation: PresentationType, source: Slide) -> Slide:
    """Clone a slide inside the same OPC package and remap relationship IDs."""
    destination = presentation.slides.add_slide(source.slide_layout)
    _remove_layout_placeholders(destination)
    relationship_map = _clone_shape_relationships(source, destination)

    destination_tree = destination.shapes._spTree  # pyright: ignore[reportPrivateUsage]
    for shape in source.shapes:
        element = deepcopy(shape.element)
        _remap_relationship_ids(element, relationship_map)
        destination_tree.insert_element_before(element, "p:extLst")
    return destination


def _remove_layout_placeholders(slide: Slide) -> None:
    tree = slide.shapes._spTree  # pyright: ignore[reportPrivateUsage]
    for shape in list(slide.shapes):
        tree.remove(shape.element)


def _clone_shape_relationships(source: Slide, destination: Slide) -> dict[str, str]:
    relationship_map: dict[str, str] = {}
    source_relationships: Any = source.part.rels
    destination_relationships: Any = destination.part.rels
    for relationship in source_relationships.values():
        if relationship.reltype in {RT.SLIDE_LAYOUT, RT.NOTES_SLIDE}:
            continue
        target = (
            relationship.target_ref
            if relationship.is_external
            else relationship.target_part
        )
        new_id = destination_relationships._add_relationship(
            relationship.reltype,
            target,
            is_external=relationship.is_external,
        )
        relationship_map[relationship.rId] = new_id
    return relationship_map


def _remap_relationship_ids(element: Any, relationship_map: dict[str, str]) -> None:
    for descendant in element.iter():
        for attribute, value in tuple(descendant.attrib.items()):
            if value in relationship_map:
                descendant.set(attribute, relationship_map[value])
